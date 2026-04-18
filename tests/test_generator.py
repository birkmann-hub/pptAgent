"""Tests for the PresentationGenerator (end-to-end .pptx creation)."""

import pytest
from datetime import date
from pathlib import Path
import sys
import tempfile

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from pptAgent.models import SectorType, ProjectType, SlideType, SlideContent
from pptAgent.agent import PresentationAgent
from pptAgent.generator import PresentationGenerator
from tests.test_models import make_request


def _generate(duration: int = 60, **kwargs) -> Path:
    """Helper: run the full pipeline and return the output path."""
    with tempfile.TemporaryDirectory() as tmp:
        req = make_request(
            duration_minutes=duration,
            output_path=str(Path(tmp) / "test.pptx"),
            **kwargs,
        )
        agent = PresentationAgent()
        plan = agent.run(req)
        gen = PresentationGenerator()
        return gen.generate(plan)


class TestPresentationGenerator:
    def test_generates_pptx_file(self):
        with tempfile.TemporaryDirectory() as tmp:
            req = make_request(
                duration_minutes=60,
                output_path=str(Path(tmp) / "out.pptx"),
            )
            agent = PresentationAgent()
            plan = agent.run(req)
            gen = PresentationGenerator()
            path = gen.generate(plan)
            assert path.exists()
            assert path.suffix == ".pptx"

    def test_pptx_is_non_empty(self):
        with tempfile.TemporaryDirectory() as tmp:
            req = make_request(
                duration_minutes=60,
                output_path=str(Path(tmp) / "out.pptx"),
            )
            agent = PresentationAgent()
            plan = agent.run(req)
            gen = PresentationGenerator()
            path = gen.generate(plan)
            assert path.stat().st_size > 5000

    def test_generates_correct_slide_count(self):
        """The generator should create exactly as many slides as the plan."""
        from pptx import Presentation as PptxPresentation

        with tempfile.TemporaryDirectory() as tmp:
            req = make_request(
                duration_minutes=90,
                output_path=str(Path(tmp) / "out.pptx"),
            )
            agent = PresentationAgent()
            plan = agent.run(req)
            gen = PresentationGenerator()
            path = gen.generate(plan)

            prs = PptxPresentation(str(path))
            assert len(prs.slides) == len(plan.slides)

    def test_short_meeting(self):
        with tempfile.TemporaryDirectory() as tmp:
            req = make_request(
                duration_minutes=15,
                output_path=str(Path(tmp) / "short.pptx"),
            )
            agent = PresentationAgent()
            plan = agent.run(req)
            gen = PresentationGenerator()
            path = gen.generate(plan)
            assert path.exists()

    def test_long_meeting(self):
        with tempfile.TemporaryDirectory() as tmp:
            req = make_request(
                duration_minutes=180,
                output_path=str(Path(tmp) / "long.pptx"),
            )
            agent = PresentationAgent()
            plan = agent.run(req)
            gen = PresentationGenerator()
            path = gen.generate(plan)
            assert path.exists()

    def test_creates_parent_directory(self):
        with tempfile.TemporaryDirectory() as tmp:
            req = make_request(
                duration_minutes=30,
                output_path=str(Path(tmp) / "subdir" / "nested" / "out.pptx"),
            )
            agent = PresentationAgent()
            plan = agent.run(req)
            gen = PresentationGenerator()
            path = gen.generate(plan)
            assert path.exists()

    def test_all_slide_types_render_without_error(self):
        """Each slide type should render without raising an exception."""
        from pptAgent.slide_builder import SLIDE_BUILDERS, build_slide
        from pptAgent.agent import _build_fallback_content
        from pptx import Presentation as PptxPresentation

        req = make_request(duration_minutes=120)
        prs = PptxPresentation()

        for i, slide_type in enumerate(SlideType, start=1):
            content = _build_fallback_content(slide_type, req, i)
            build_slide(prs, content, i)  # must not raise

        assert len(prs.slides) == len(SlideType)

    @pytest.mark.parametrize("sector", list(SectorType))
    def test_sector_variants(self, sector, tmp_path):
        req = make_request(
            sector=sector,
            output_path=str(tmp_path / "out.pptx"),
            duration_minutes=30,
        )
        agent = PresentationAgent()
        plan = agent.run(req)
        gen = PresentationGenerator()
        path = gen.generate(plan)
        assert path.exists()

    @pytest.mark.parametrize("project_type", list(ProjectType))
    def test_project_type_variants(self, project_type, tmp_path):
        req = make_request(
            project_type=project_type,
            output_path=str(tmp_path / "out.pptx"),
            duration_minutes=30,
        )
        agent = PresentationAgent()
        plan = agent.run(req)
        gen = PresentationGenerator()
        path = gen.generate(plan)
        assert path.exists()
