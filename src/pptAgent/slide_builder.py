"""Renders individual slides onto a python-pptx Presentation object.

Each public ``build_*`` function accepts a ``pptx.Presentation`` instance and
a ``SlideContent`` object and appends one fully-styled slide.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .models import SlideContent, SlideType
from .corporate_design import (
    DARK_BLUE, MEDIUM_BLUE, LIGHT_BLUE, ACCENT_ORANGE,
    WHITE, LIGHT_GRAY, DARK_TEXT, GRAY, MID_GRAY,
    FONT_FAMILY, FONT_FAMILY_HEADINGS,
    SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU,
    HEADER_HEIGHT, CONTENT_TOP, CONTENT_HEIGHT,
    FOOTER_TOP, FOOTER_HEIGHT,
    LEFT_MARGIN, CONTENT_WIDTH,
    resolve_color,
)

# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation) -> object:
    """Add a blank slide using the 'Blank' layout (index 6)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _add_rect(slide, left, top, width, height, fill_color: RGBColor | None = None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()  # no border
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def _add_textbox(
    slide,
    left, top, width, height,
    text: str,
    font_size: int = 14,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = DARK_TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_family: str = FONT_FAMILY,
    word_wrap: bool = True,
) -> object:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_family
    return txBox


def _add_bullet_list(
    slide,
    left, top, width, height,
    items: list[str],
    font_size: int = 14,
    color: RGBColor = DARK_TEXT,
    numbered: bool = False,
    indent_pt: int = 12,
) -> object:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        prefix = f"{i + 1}.  " if numbered else "▪  "
        run.text = f"{prefix}{item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = FONT_FAMILY
    return txBox


def _add_colored_label(slide, left, top, width, height, text, font_size=12, bg_color=DARK_BLUE, text_color=WHITE):
    """A small coloured badge / label."""
    rect = _add_rect(slide, left, top, width, height, bg_color)
    tf = rect.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = FONT_FAMILY
    return rect


# ---------------------------------------------------------------------------
# Shared chrome: header band + footer band + slide number
# ---------------------------------------------------------------------------

def _add_chrome(slide, title_text: str, slide_number: int | None = None):
    """Adds the dark-blue header band with slide title and footer line."""
    # Header band
    header_rect = _add_rect(
        slide,
        left=Emu(0), top=Emu(0),
        width=Emu(SLIDE_WIDTH_EMU), height=HEADER_HEIGHT,
        fill_color=DARK_BLUE,
    )
    # Title text in header
    _add_textbox(
        slide,
        left=LEFT_MARGIN, top=Emu(int(SLIDE_HEIGHT_EMU * 0.02)),
        width=CONTENT_WIDTH, height=HEADER_HEIGHT,
        text=title_text,
        font_size=18,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.LEFT,
        font_family=FONT_FAMILY_HEADINGS,
    )
    # Footer accent line
    _add_rect(
        slide,
        left=Emu(0), top=FOOTER_TOP,
        width=Emu(SLIDE_WIDTH_EMU), height=Emu(int(SLIDE_HEIGHT_EMU * 0.004)),
        fill_color=ACCENT_ORANGE,
    )
    # Slide number (optional)
    if slide_number is not None:
        _add_textbox(
            slide,
            left=Emu(SLIDE_WIDTH_EMU - int(SLIDE_WIDTH_EMU * 0.06)),
            top=FOOTER_TOP,
            width=Emu(int(SLIDE_WIDTH_EMU * 0.05)),
            height=FOOTER_HEIGHT,
            text=str(slide_number),
            font_size=10,
            color=GRAY,
            align=PP_ALIGN.RIGHT,
        )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_cover(prs: Presentation, content: SlideContent, slide_number: int = 1) -> None:
    """Full-bleed dark cover slide."""
    slide = _blank_slide(prs)

    # Dark-blue background
    _add_rect(slide, Emu(0), Emu(0), Emu(SLIDE_WIDTH_EMU), Emu(SLIDE_HEIGHT_EMU), DARK_BLUE)

    # Accent stripe on left edge
    _add_rect(
        slide,
        left=Emu(0), top=Emu(0),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.01)),
        height=Emu(SLIDE_HEIGHT_EMU),
        fill_color=ACCENT_ORANGE,
    )

    # Main title
    _add_textbox(
        slide,
        left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.25)),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.88)),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.25)),
        text=content.title or "Presentation Title",
        font_size=36,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.LEFT,
        font_family=FONT_FAMILY_HEADINGS,
    )
    # Subtitle
    _add_textbox(
        slide,
        left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.52)),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.88)),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.12)),
        text=content.subtitle or "",
        font_size=20,
        color=LIGHT_BLUE,
        align=PP_ALIGN.LEFT,
    )
    # Meta line (sector | project type | date)
    _add_textbox(
        slide,
        left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.66)),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.88)),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.10)),
        text=content.body_text or "",
        font_size=13,
        color=WHITE,
        align=PP_ALIGN.LEFT,
    )
    # Footer – presenter name
    _add_textbox(
        slide,
        left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.82)),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.88)),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.10)),
        text=content.footnote or "",
        font_size=12,
        color=LIGHT_BLUE,
        align=PP_ALIGN.LEFT,
    )


def build_agenda(prs: Presentation, content: SlideContent, slide_number: int = 2) -> None:
    slide = _blank_slide(prs)
    _add_chrome(slide, content.title or "Agenda", slide_number)

    items = content.bullet_points or []
    item_height = Emu(int(CONTENT_HEIGHT / max(len(items), 1)))
    top = CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.02))

    for i, item in enumerate(items):
        # Numbered badge
        badge_w = Emu(int(SLIDE_WIDTH_EMU * 0.05))
        _add_colored_label(
            slide,
            left=LEFT_MARGIN,
            top=top + i * item_height,
            width=badge_w,
            height=Emu(int(item_height * 0.7)),
            text=f"{i + 1:02d}",
            font_size=14,
            bg_color=DARK_BLUE,
            text_color=WHITE,
        )
        # Item text
        _add_textbox(
            slide,
            left=LEFT_MARGIN + badge_w + Emu(int(SLIDE_WIDTH_EMU * 0.01)),
            top=top + i * item_height,
            width=CONTENT_WIDTH - badge_w - Emu(int(SLIDE_WIDTH_EMU * 0.01)),
            height=Emu(int(item_height * 0.7)),
            text=item,
            font_size=16,
            color=DARK_TEXT,
            align=PP_ALIGN.LEFT,
        )

    if content.footnote:
        _add_textbox(
            slide,
            left=LEFT_MARGIN,
            top=FOOTER_TOP - Emu(int(SLIDE_HEIGHT_EMU * 0.03)),
            width=CONTENT_WIDTH,
            height=FOOTER_HEIGHT,
            text=content.footnote,
            font_size=10,
            color=GRAY,
        )


def build_executive_summary(prs: Presentation, content: SlideContent, slide_number: int = 3) -> None:
    slide = _blank_slide(prs)
    _add_chrome(slide, content.title or "Executive Summary", slide_number)

    # Bullet points
    list_height = CONTENT_HEIGHT - Emu(int(SLIDE_HEIGHT_EMU * 0.18))
    _add_bullet_list(
        slide,
        left=LEFT_MARGIN,
        top=CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.02)),
        width=CONTENT_WIDTH,
        height=list_height,
        items=content.bullet_points,
        font_size=15,
        color=DARK_TEXT,
    )

    # Highlight / bottom-line box
    if content.highlight_box:
        box_top = FOOTER_TOP - Emu(int(SLIDE_HEIGHT_EMU * 0.12))
        rect = _add_rect(
            slide,
            left=LEFT_MARGIN, top=box_top,
            width=CONTENT_WIDTH,
            height=Emu(int(SLIDE_HEIGHT_EMU * 0.10)),
            fill_color=DARK_BLUE,
        )
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = content.highlight_box
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT_FAMILY


def build_context_background(prs: Presentation, content: SlideContent, slide_number: int = 4) -> None:
    slide = _blank_slide(prs)
    _add_chrome(slide, content.title or "Context & Background", slide_number)

    col_width = Emu(int(CONTENT_WIDTH / 2 - SLIDE_WIDTH_EMU * 0.01))
    col_top = CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.02))
    col_height = CONTENT_HEIGHT - Emu(int(SLIDE_HEIGHT_EMU * 0.06))

    # Divider line
    _add_rect(
        slide,
        left=LEFT_MARGIN + col_width + Emu(int(SLIDE_WIDTH_EMU * 0.005)),
        top=col_top,
        width=Emu(int(SLIDE_WIDTH_EMU * 0.002)),
        height=col_height,
        fill_color=LIGHT_GRAY,
    )

    for col_idx, (heading, bullets) in enumerate([
        (content.left_heading or "Industry Trends", content.left_bullets),
        (content.right_heading or "Company Background", content.right_bullets),
    ]):
        col_left = LEFT_MARGIN + col_idx * (col_width + Emu(int(SLIDE_WIDTH_EMU * 0.02)))
        # Heading
        _add_textbox(
            slide,
            left=col_left, top=col_top,
            width=col_width,
            height=Emu(int(SLIDE_HEIGHT_EMU * 0.07)),
            text=heading,
            font_size=16,
            bold=True,
            color=DARK_BLUE,
        )
        # Bullets
        _add_bullet_list(
            slide,
            left=col_left,
            top=col_top + Emu(int(SLIDE_HEIGHT_EMU * 0.08)),
            width=col_width,
            height=col_height - Emu(int(SLIDE_HEIGHT_EMU * 0.08)),
            items=bullets,
            font_size=13,
            color=DARK_TEXT,
        )

    if content.footnote:
        _add_textbox(
            slide, LEFT_MARGIN,
            FOOTER_TOP - Emu(int(SLIDE_HEIGHT_EMU * 0.02)),
            CONTENT_WIDTH, FOOTER_HEIGHT,
            content.footnote, font_size=10, color=GRAY,
        )


def build_problem_statement(prs: Presentation, content: SlideContent, slide_number: int = 5) -> None:
    slide = _blank_slide(prs)
    _add_chrome(slide, content.title or "Problem Statement", slide_number)

    top = CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.02))

    # Core question box (accent colour)
    if content.highlight_box:
        box_height = Emu(int(SLIDE_HEIGHT_EMU * 0.12))
        rect = _add_rect(slide, LEFT_MARGIN, top, CONTENT_WIDTH, box_height, ACCENT_ORANGE)
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = content.highlight_box
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT_FAMILY
        top = top + box_height + Emu(int(SLIDE_HEIGHT_EMU * 0.04))

    # Problem drivers
    _add_bullet_list(
        slide,
        left=LEFT_MARGIN, top=top,
        width=CONTENT_WIDTH,
        height=CONTENT_HEIGHT - Emu(int(SLIDE_HEIGHT_EMU * 0.20)),
        items=content.bullet_points,
        font_size=14,
        color=DARK_TEXT,
    )

    if content.body_text:
        _add_textbox(
            slide, LEFT_MARGIN,
            FOOTER_TOP - Emu(int(SLIDE_HEIGHT_EMU * 0.06)),
            CONTENT_WIDTH,
            Emu(int(SLIDE_HEIGHT_EMU * 0.05)),
            content.body_text, font_size=13, italic=True, color=GRAY,
        )


def build_methodology(prs: Presentation, content: SlideContent, slide_number: int = 6) -> None:
    slide = _blank_slide(prs)
    _add_chrome(slide, content.title or "Our Approach", slide_number)

    steps = content.process_steps or []
    if not steps:
        return

    top = CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.06))
    box_height = Emu(int(SLIDE_HEIGHT_EMU * 0.52))
    arrow_width = Emu(int(SLIDE_WIDTH_EMU * 0.02))

    n = len(steps)
    available_width = CONTENT_WIDTH - arrow_width * (n - 1)
    box_width = Emu(int(available_width / n))

    for i, step in enumerate(steps):
        box_left = LEFT_MARGIN + i * (box_width + arrow_width)
        is_active = step.get("active", False)
        bg = ACCENT_ORANGE if is_active else DARK_BLUE

        # Phase box
        rect = _add_rect(slide, box_left, top, box_width, box_height, bg)
        tf = rect.text_frame
        tf.word_wrap = True

        # Phase number
        p_num = tf.paragraphs[0]
        p_num.alignment = PP_ALIGN.CENTER
        r = p_num.add_run()
        r.text = step.get("label", f"Phase {i + 1}")
        r.font.size = Pt(11)
        r.font.color.rgb = WHITE
        r.font.name = FONT_FAMILY

        # Phase heading
        p_head = tf.add_paragraph()
        p_head.alignment = PP_ALIGN.CENTER
        p_head.space_before = Pt(6)
        rh = p_head.add_run()
        rh.text = step.get("heading", "")
        rh.font.size = Pt(14)
        rh.font.bold = True
        rh.font.color.rgb = WHITE
        rh.font.name = FONT_FAMILY_HEADINGS

        # Description
        p_desc = tf.add_paragraph()
        p_desc.alignment = PP_ALIGN.CENTER
        p_desc.space_before = Pt(8)
        rd = p_desc.add_run()
        rd.text = step.get("description", "")
        rd.font.size = Pt(11)
        rd.font.color.rgb = LIGHT_GRAY
        rd.font.name = FONT_FAMILY

        # Arrow connector (except after last box)
        if i < n - 1:
            arrow_left = box_left + box_width
            _add_textbox(
                slide,
                left=arrow_left,
                top=top + Emu(int(box_height / 2) - int(SLIDE_HEIGHT_EMU * 0.03)),
                width=arrow_width,
                height=Emu(int(SLIDE_HEIGHT_EMU * 0.06)),
                text="▶",
                font_size=16,
                color=MID_GRAY,
                align=PP_ALIGN.CENTER,
            )

    if content.footnote:
        _add_textbox(
            slide, LEFT_MARGIN,
            FOOTER_TOP - Emu(int(SLIDE_HEIGHT_EMU * 0.02)),
            CONTENT_WIDTH, FOOTER_HEIGHT,
            content.footnote, font_size=10, color=GRAY,
        )


def build_findings(prs: Presentation, content: SlideContent, slide_number: int = 7) -> None:
    slide = _blank_slide(prs)
    _add_chrome(slide, content.title or "Key Findings", slide_number)

    top = CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.02))
    col_width = Emu(int(CONTENT_WIDTH / 2 - SLIDE_WIDTH_EMU * 0.01))

    # Finding badge (left)
    badge_h = Emu(int(SLIDE_HEIGHT_EMU * 0.06))
    if content.section_number:
        _add_colored_label(
            slide,
            left=LEFT_MARGIN, top=top,
            width=Emu(int(SLIDE_WIDTH_EMU * 0.10)),
            height=badge_h,
            text=f"Finding {content.section_number}",
            font_size=11,
            bg_color=ACCENT_ORANGE,
            text_color=WHITE,
        )
        top = top + badge_h + Emu(int(SLIDE_HEIGHT_EMU * 0.01))

    # Finding headline (left column)
    _add_textbox(
        slide,
        left=LEFT_MARGIN, top=top,
        width=col_width,
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.12)),
        text=content.highlight_box or content.subtitle or "",
        font_size=18,
        bold=True,
        color=DARK_BLUE,
    )
    # Supporting bullets (left column)
    _add_bullet_list(
        slide,
        left=LEFT_MARGIN,
        top=top + Emu(int(SLIDE_HEIGHT_EMU * 0.14)),
        width=col_width,
        height=CONTENT_HEIGHT - Emu(int(SLIDE_HEIGHT_EMU * 0.22)),
        items=content.bullet_points,
        font_size=13,
        color=DARK_TEXT,
    )

    # Right column: chart placeholder
    right_left = LEFT_MARGIN + col_width + Emu(int(SLIDE_WIDTH_EMU * 0.02))
    placeholder_rect = _add_rect(
        slide,
        left=right_left,
        top=CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.02)),
        width=col_width,
        height=CONTENT_HEIGHT - Emu(int(SLIDE_HEIGHT_EMU * 0.04)),
        fill_color=LIGHT_GRAY,
    )
    tf = placeholder_rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = content.right_heading or "[Chart / Visual]"
    r.font.size = Pt(13)
    r.font.color.rgb = GRAY
    r.font.name = FONT_FAMILY


def build_recommendations(prs: Presentation, content: SlideContent, slide_number: int = 8) -> None:
    slide = _blank_slide(prs)
    _add_chrome(slide, content.title or "Our Recommendations", slide_number)

    cards = content.recommendation_cards or []
    if not cards:
        _add_bullet_list(
            slide,
            left=LEFT_MARGIN,
            top=CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.02)),
            width=CONTENT_WIDTH,
            height=CONTENT_HEIGHT,
            items=content.bullet_points,
            font_size=14,
            color=DARK_TEXT,
        )
        return

    n = len(cards)
    top = CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.04))
    card_height = Emu(int((CONTENT_HEIGHT - SLIDE_HEIGHT_EMU * 0.08) / n))
    card_gap = Emu(int(SLIDE_HEIGHT_EMU * 0.01))

    for i, card in enumerate(cards):
        card_top = top + i * (card_height + card_gap)

        # Priority badge
        badge_w = Emu(int(SLIDE_WIDTH_EMU * 0.06))
        _add_colored_label(
            slide,
            left=LEFT_MARGIN, top=card_top,
            width=badge_w, height=card_height,
            text=card.get("priority", f"{i + 1:02d}"),
            font_size=20,
            bg_color=DARK_BLUE,
            text_color=WHITE,
        )

        # Card body
        text_left = LEFT_MARGIN + badge_w + Emu(int(SLIDE_WIDTH_EMU * 0.01))
        text_width = CONTENT_WIDTH - badge_w - Emu(int(SLIDE_WIDTH_EMU * 0.20))
        rect = _add_rect(slide, text_left, card_top, text_width, card_height, LIGHT_GRAY)
        tf = rect.text_frame
        tf.word_wrap = True
        p_head = tf.paragraphs[0]
        r_head = p_head.add_run()
        r_head.text = card.get("heading", "")
        r_head.font.size = Pt(14)
        r_head.font.bold = True
        r_head.font.color.rgb = DARK_BLUE
        r_head.font.name = FONT_FAMILY_HEADINGS

        p_desc = tf.add_paragraph()
        r_desc = p_desc.add_run()
        r_desc.text = card.get("description", "")
        r_desc.font.size = Pt(12)
        r_desc.font.color.rgb = DARK_TEXT
        r_desc.font.name = FONT_FAMILY

        # Impact / effort tags on right
        tag_left = text_left + text_width + Emu(int(SLIDE_WIDTH_EMU * 0.01))
        tag_width = Emu(int(SLIDE_WIDTH_EMU * 0.12))
        _add_colored_label(
            slide,
            left=tag_left, top=card_top,
            width=tag_width,
            height=Emu(int(card_height / 2 - SLIDE_HEIGHT_EMU * 0.005)),
            text=f"Impact: {card.get('impact', 'High')}",
            font_size=10,
            bg_color=MEDIUM_BLUE,
            text_color=WHITE,
        )
        _add_colored_label(
            slide,
            left=tag_left,
            top=card_top + Emu(int(card_height / 2 + SLIDE_HEIGHT_EMU * 0.005)),
            width=tag_width,
            height=Emu(int(card_height / 2 - SLIDE_HEIGHT_EMU * 0.005)),
            text=f"Effort: {card.get('effort', 'Medium')}",
            font_size=10,
            bg_color=ACCENT_ORANGE,
            text_color=WHITE,
        )

    if content.footnote:
        _add_textbox(
            slide, LEFT_MARGIN,
            FOOTER_TOP - Emu(int(SLIDE_HEIGHT_EMU * 0.02)),
            CONTENT_WIDTH, FOOTER_HEIGHT,
            content.footnote, font_size=10, color=GRAY,
        )


def build_roadmap(prs: Presentation, content: SlideContent, slide_number: int = 9) -> None:
    slide = _blank_slide(prs)
    _add_chrome(slide, content.title or "Implementation Roadmap", slide_number)

    columns = content.timeline_columns or []
    lanes = content.swim_lanes or []

    if not columns or not lanes:
        _add_textbox(
            slide, LEFT_MARGIN,
            CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.15)),
            CONTENT_WIDTH, CONTENT_HEIGHT,
            "Roadmap – to be detailed in project planning",
            font_size=14, color=GRAY, align=PP_ALIGN.CENTER,
        )
        return

    header_top = CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.02))
    header_height = Emu(int(SLIDE_HEIGHT_EMU * 0.08))
    row_height = Emu(int((CONTENT_HEIGHT - SLIDE_HEIGHT_EMU * 0.14) / max(len(lanes), 1)))
    label_width = Emu(int(SLIDE_WIDTH_EMU * 0.15))
    timeline_width = CONTENT_WIDTH - label_width
    col_width = Emu(int(timeline_width / len(columns)))

    # Column headers
    for j, col_label in enumerate(columns):
        col_left = LEFT_MARGIN + label_width + j * col_width
        hdr_rect = _add_rect(slide, col_left, header_top, col_width - Emu(int(SLIDE_WIDTH_EMU * 0.002)), header_height, DARK_BLUE)
        tf = hdr_rect.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = col_label
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT_FAMILY

    # Swim lanes
    for i, lane in enumerate(lanes):
        row_top = header_top + header_height + i * row_height
        # Workstream label
        label_rect = _add_rect(
            slide, LEFT_MARGIN, row_top,
            label_width - Emu(int(SLIDE_WIDTH_EMU * 0.01)), row_height - Emu(int(SLIDE_HEIGHT_EMU * 0.01)),
            LIGHT_GRAY,
        )
        tf = label_rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = lane.get("workstream", f"Workstream {i + 1}")
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = DARK_BLUE
        r.font.name = FONT_FAMILY

        # Bars
        for bar in lane.get("bars", []):
            start_col = max(0, bar.get("start", 1) - 1)
            span = max(1, bar.get("span", 1))
            bar_left = LEFT_MARGIN + label_width + start_col * col_width
            bar_width = span * col_width - Emu(int(SLIDE_WIDTH_EMU * 0.005))
            bar_rect = _add_rect(
                slide, bar_left,
                row_top + Emu(int(row_height * 0.15)),
                bar_width,
                Emu(int(row_height * 0.7)),
                MEDIUM_BLUE,
            )
            tf = bar_rect.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = bar.get("label", "")
            r.font.size = Pt(10)
            r.font.color.rgb = WHITE
            r.font.name = FONT_FAMILY


def build_next_steps(prs: Presentation, content: SlideContent, slide_number: int = 10) -> None:
    slide = _blank_slide(prs)
    _add_chrome(slide, content.title or "Next Steps", slide_number)

    rows = content.table_rows or []
    if not rows:
        return

    headers = ["#", "Action Item", "Owner", "Due Date", "Status"]
    col_widths_frac = [0.05, 0.45, 0.18, 0.15, 0.13]
    header_height = Emu(int(SLIDE_HEIGHT_EMU * 0.07))
    row_height = Emu(int((CONTENT_HEIGHT - SLIDE_HEIGHT_EMU * 0.10) / (len(rows) + 1)))
    table_top = CONTENT_TOP + Emu(int(SLIDE_HEIGHT_EMU * 0.03))

    col_widths = [Emu(int(CONTENT_WIDTH * f)) for f in col_widths_frac]

    # Header row
    x = LEFT_MARGIN
    for j, (header, cw) in enumerate(zip(headers, col_widths)):
        rect = _add_rect(slide, x, table_top, cw - Emu(int(SLIDE_WIDTH_EMU * 0.002)), header_height, DARK_BLUE)
        tf = rect.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = header
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT_FAMILY
        x = x + cw

    # Data rows
    for i, row_data in enumerate(rows):
        row_top = table_top + header_height + i * row_height
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        x = LEFT_MARGIN
        for j, (cell_text, cw) in enumerate(zip(row_data, col_widths)):
            rect = _add_rect(slide, x, row_top, cw - Emu(int(SLIDE_WIDTH_EMU * 0.002)), row_height - Emu(int(SLIDE_HEIGHT_EMU * 0.005)), bg)
            tf = rect.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j in (0, 2, 3, 4) else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(cell_text)
            r.font.size = Pt(11)
            r.font.color.rgb = DARK_TEXT
            r.font.name = FONT_FAMILY
            x = x + cw


def build_section_divider(prs: Presentation, content: SlideContent, slide_number: int | None = None) -> None:
    slide = _blank_slide(prs)

    # Full dark-blue background
    _add_rect(slide, Emu(0), Emu(0), Emu(SLIDE_WIDTH_EMU), Emu(SLIDE_HEIGHT_EMU), DARK_BLUE)

    # Horizontal accent stripe
    _add_rect(
        slide,
        left=Emu(0),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.40)),
        width=Emu(SLIDE_WIDTH_EMU),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.004)),
        fill_color=ACCENT_ORANGE,
    )

    # Section number (large, faded)
    _add_textbox(
        slide,
        left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.10)),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.20)),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.35)),
        text=content.section_number or "",
        font_size=80,
        bold=True,
        color=ACCENT_ORANGE,
        align=PP_ALIGN.LEFT,
        font_family=FONT_FAMILY_HEADINGS,
    )
    # Section title
    _add_textbox(
        slide,
        left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.44)),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.88)),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.20)),
        text=content.title or "",
        font_size=36,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.LEFT,
        font_family=FONT_FAMILY_HEADINGS,
    )
    # Section description
    _add_textbox(
        slide,
        left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.65)),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.88)),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.15)),
        text=content.subtitle or "",
        font_size=16,
        color=LIGHT_BLUE,
        align=PP_ALIGN.LEFT,
    )


def build_closing(prs: Presentation, content: SlideContent, slide_number: int | None = None) -> None:
    slide = _blank_slide(prs)

    # Full dark-blue background
    _add_rect(slide, Emu(0), Emu(0), Emu(SLIDE_WIDTH_EMU), Emu(SLIDE_HEIGHT_EMU), DARK_BLUE)

    # Left accent stripe
    _add_rect(
        slide,
        left=Emu(0), top=Emu(0),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.01)),
        height=Emu(SLIDE_HEIGHT_EMU),
        fill_color=ACCENT_ORANGE,
    )

    # "Thank You" or custom main text
    _add_textbox(
        slide,
        left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.25)),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.88)),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.22)),
        text=content.title or "Thank You",
        font_size=48,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.LEFT,
        font_family=FONT_FAMILY_HEADINGS,
    )
    # Tagline
    _add_textbox(
        slide,
        left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.50)),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.88)),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.12)),
        text=content.subtitle or "We look forward to your questions and next steps.",
        font_size=18,
        color=LIGHT_BLUE,
        align=PP_ALIGN.LEFT,
    )
    # Contact info
    _add_textbox(
        slide,
        left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
        top=Emu(int(SLIDE_HEIGHT_EMU * 0.70)),
        width=Emu(int(SLIDE_WIDTH_EMU * 0.88)),
        height=Emu(int(SLIDE_HEIGHT_EMU * 0.10)),
        text=content.body_text or "",
        font_size=14,
        color=WHITE,
        align=PP_ALIGN.LEFT,
    )
    if content.footnote:
        _add_textbox(
            slide,
            left=Emu(int(SLIDE_WIDTH_EMU * 0.06)),
            top=Emu(int(SLIDE_HEIGHT_EMU * 0.82)),
            width=Emu(int(SLIDE_WIDTH_EMU * 0.88)),
            height=Emu(int(SLIDE_HEIGHT_EMU * 0.08)),
            text=content.footnote,
            font_size=12,
            color=LIGHT_BLUE,
            align=PP_ALIGN.LEFT,
        )


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

SLIDE_BUILDERS = {
    SlideType.COVER: build_cover,
    SlideType.AGENDA: build_agenda,
    SlideType.EXECUTIVE_SUMMARY: build_executive_summary,
    SlideType.CONTEXT_BACKGROUND: build_context_background,
    SlideType.PROBLEM_STATEMENT: build_problem_statement,
    SlideType.METHODOLOGY: build_methodology,
    SlideType.FINDINGS: build_findings,
    SlideType.RECOMMENDATIONS: build_recommendations,
    SlideType.ROADMAP: build_roadmap,
    SlideType.NEXT_STEPS: build_next_steps,
    SlideType.SECTION_DIVIDER: build_section_divider,
    SlideType.CLOSING: build_closing,
}


def build_slide(prs: Presentation, content: SlideContent, slide_number: int) -> None:
    """Dispatch to the correct builder for the given slide type."""
    builder = SLIDE_BUILDERS.get(content.slide_type)
    if builder is None:
        raise ValueError(f"No builder registered for slide type: {content.slide_type}")
    builder(prs, content, slide_number)
